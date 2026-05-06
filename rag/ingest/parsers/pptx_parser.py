from __future__ import annotations

from .base import BaseParser, ParseResult, TextChunk


class PptxParser(BaseParser):

    def can_parse(self, file_path: str) -> bool:
        return self._ext(file_path) == ".pptx"

    def parse(self, file_path: str) -> ParseResult:
        try:
            from pptx import Presentation
        except ImportError:
            return ParseResult(
                source_path=file_path, doc_type="pptx",
                errors=["python-pptx가 설치되지 않았습니다."]
            )

        chunks: list[TextChunk] = []
        errors: list[str] = []
        title = ""

        try:
            prs = Presentation(file_path)
            try:
                title = prs.core_properties.title or ""
            except Exception:
                title = ""

            for slide_idx, slide in enumerate(prs.slides, start=1):
                texts: list[str] = []
                slide_title = ""

                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for para in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in para.runs).strip()
                        if line:
                            texts.append(line)
                            if not slide_title and shape == slide.shapes.title:
                                slide_title = line

                if hasattr(slide, "has_notes_slide") and slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        texts.append(f"[발표자 노트]\n{notes}")

                if not texts:
                    continue

                content = f"[슬라이드 {slide_idx}" + (f": {slide_title}" if slide_title else "") + "]\n"
                content += "\n".join(texts)

                chunks.append(TextChunk(
                    content=content,
                    chunk_index=slide_idx - 1,
                    page_or_line=slide_idx,
                    metadata={"slide_number": slide_idx, "slide_title": slide_title},
                ))
        except Exception as e:
            errors.append(str(e))

        return ParseResult(
            source_path=file_path,
            doc_type="pptx",
            chunks=chunks,
            metadata={"title": title, "slide_count": len(chunks)},
            errors=errors,
        )
